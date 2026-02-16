from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_title_slide(prs, title_text, subtitle_text, presenters):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(0, 0, 0))

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255) # Cyan
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    
    tf = subtitle.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = subtitle_text
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "\nPresented By:"
    p.font.color.rgb = RGBColor(255, 0, 255) # Magenta
    p.font.bold = True
    p.font.size = Pt(20)

    for presenter in presenters:
        p = tf.add_paragraph()
        p.text = presenter
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(18)

def add_content_slide(prs, title_text, content_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(0, 0, 0))

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 255) # Cyan
    title.text_frame.paragraphs[0].font.bold = True

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(22)
        p.space_after = Pt(14)
        
        # Check for sub-points (simple implementation: strings starting with '  -')
        if point.strip().startswith("-"):
             p.level = 1
             p.text = point.strip()[1:].strip()

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(prs, 
                   "Online Coding Platform", 
                   "College: Birla Vishvakarma Mahavidhyalaya\nMentor: Prof. Nilesh Prajapati",
                   ["Jagdish Pathakji", "Yug Dholakiya", "Rudra Gohil"])

    # Slide 2: Aim of the Project
    add_content_slide(prs, "Aim of the Project", [
        "Develop an Online Coding Platform with AI-assisted guidance.",
        "Enhance students' problem-solving and interview preparation skills.",
        "Provide a unified platform for:",
        "  - DSA practice",
        "  - AI-based coding assistance",
        "  - 1-vs-1 coding challenges",
        "  - Interview simulations"
    ])

    # Slide 3: Project Scope
    add_content_slide(prs, "Project Scope", [
        "Web-based system designed for technical interview prep.",
        "Target Audience: Students & individuals improving coding skills.",
        "Key Features:",
        "  - Practice Data Structures and Algorithms (DSA)",
        "  - Real-time AI guidance",
        "  - Competitive 1-vs-1 challenges",
        "  - Mock Technical Interviews"
    ])

    # Slide 4: Project Objectives
    add_content_slide(prs, "Project Objectives", [
        "Solve Coding Problems: Practice across various difficulty levels.",
        "AI-Assisted Support: Context-aware hints and optimization tips.",
        "Interview Preparation: AI-driven and live simulation tools.",
        "User Engagement: Interactive and competitive environment."
    ])

    # Slide 5: Project Modules Overview
    add_content_slide(prs, "Project Modules", [
        "1. Authentication Module",
        "2. Problem Management Module",
        "3. Interview Practice Module",
        "4. One-on-One Coding Challenge Module",
        "5. Progress Tracking Module",
        "6. Profile Section Module"
    ])

    # Slide 6: System Requirements
    add_content_slide(prs, "System Requirements", [
        "Hardware:",
        "  - Processor: Intel Core i5 or equivalent",
        "  - RAM: 4 GB or higher",
        "  - Internet: Stable connection",
        "Software Stack:",
        "  - Frontend: React (HTML & CSS)",
        "  - Backend: Node.js, Express.js",
        "  - Database: MongoDB",
        "  - AI & Hosting: Ollama (gpt-oss), AWS"
    ])

    # Slide 7: Detailed Integration - Core
    add_content_slide(prs, "Core Modules", [
        "Authentication:",
        "  - Secure Registration & Login",
        "  - JWT based Session Management",
        "Problem Management:",
        "  - Structured Repository (Easy/Medium/Hard)",
        "  - Hidden/Visible Test Cases",
        "  - AI Hints & Syntax Error Explanation"
    ])

    # Slide 8: Detailed Integration - Advanced
    add_content_slide(prs, "Advanced Modules", [
        "Interview Practice:",
        "  - Voice-based AI technical interviews",
        "  - Standard DSA interview rounds",
        "  - Timed sessions for realism",
        "One-on-One Challenges:",
        "  - Real-time competitive coding",
        "  - Synchronized problem sets",
        "  - Live result comparison"
    ])

    # Slide 9: Analytics & Profile
    add_content_slide(prs, "Analytics & Profile", [
        "Progress Tracking:",
        "  - Exact submission history",
        "  - Accuracy & Performance Analytics",
        "  - Visual graphs of improvement",
        "Profile Section:",
        "  - Skill Overview & Badges",
        "  - Public Profile for showcasing capability"
    ])

    # Slide 10: Feasibility Study
    add_content_slide(prs, "Feasibility Study", [
        "Technical:",
        "  - Built on robust MERN stack.",
        "  - Feasible local LLM integration (Ollama).",
        "Operational:",
        "  - User-friendly interface, no training required.",
        "  - Supports self-paced learning.",
        "Economic:",
        "  - Uses open-source technologies.",
        "  - Scalable cloud deployment (AWS)."
    ])

    # Slide 11: Team & Conclusion
    add_content_slide(prs, "Team Roles", [
        "Jagdish Pathakji: Problem Mgmt, 1v1 Challenge Module",
        "Yug Dholakiya: Authentication, Interview Practice Module",
        "Rudra Gohil: Progress Tracking, Profile Section",
        "",
        "Conclusion:",
        "The platform bridges the gap between learning and industry application through AI and competition."
    ])

    prs.save("Online_Coding_Platform_Presentation.pptx")
    print("Presentation created successfully: Online_Coding_Platform_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
